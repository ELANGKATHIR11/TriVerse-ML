"""
PowerPoint Report Generator for CodeAlpha Enterprise AI Platform.

Generates a professional slide deck using python-pptx with:
- Title slide with branding
- Dataset overview slide
- Model comparison table slide
- Metrics summary slide
- Per-model SHAP slides
- Leaderboard and recommendations finale slide
"""

from __future__ import annotations

import base64
import logging
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------

_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)      # Indigo 600
_SECONDARY = RGBColor(0x7C, 0x3A, 0xED)    # Violet 600
_ACCENT = RGBColor(0x06, 0xB6, 0xD4)       # Cyan 500
_DARK = RGBColor(0x1E, 0x1B, 0x4B)         # Indigo 950
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)        # Indigo 50
_GOLD = RGBColor(0xF5, 0x9E, 0x0B)         # Amber 500

# Slide layout indices (from blank template)
_BLANK = 6
_TITLE_CONTENT = 1


def _rgb_hex(r: int, g: int, b: int) -> str:
    return f"{r:02X}{g:02X}{b:02X}"


# ---------------------------------------------------------------------------
# Low-level PPTX helpers
# ---------------------------------------------------------------------------

def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: RGBColor,
    line_rgb: RGBColor | None = None,
    line_width: int = 0,
) -> Any:
    """Add a filled rectangle shape to a slide."""
    from pptx.util import Inches as _I
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _I(left), _I(top), _I(width), _I(height),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    if line_rgb and line_width:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = _DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> Any:
    """Add a textbox with formatted text to a slide."""
    from pptx.util import Inches as _I
    txBox = slide.shapes.add_textbox(_I(left), _I(top), _I(width), _I(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _format_val(val: Any) -> str:
    if isinstance(val, float):
        return f"{val:.4f}"
    return str(val) if val is not None else "—"


# ---------------------------------------------------------------------------
# Generator
# ---------------------------------------------------------------------------

class PPTXReportGenerator:
    """
    Builds a PowerPoint presentation for ML experiment results.

    Slide structure:
    1. Title
    2. Dataset Overview
    3. Model Comparison Table
    4. Performance Metrics Summary
    5-N. SHAP plots (one per model)
    Final. Leaderboard + Recommendations
    """

    _SLIDE_W = 10.0  # inches
    _SLIDE_H = 7.5   # inches

    def generate(
        self,
        output_path: Path,
        title: str,
        task_type: str,
        dataset_summary: dict[str, Any],
        metrics_table: list[dict[str, Any]],
        leaderboard: list[dict[str, Any]],
        shap_images: dict[str, str],
        ai_insights: str,
    ) -> Path:
        """
        Build and save the PowerPoint presentation.

        Args:
            output_path:     Destination ``.pptx`` file path.
            title:           Presentation / experiment title.
            task_type:       ML task type.
            dataset_summary: Dataset metadata dict.
            metrics_table:   Per-model evaluation metric rows.
            leaderboard:     Ranked leaderboard entries.
            shap_images:     Base64-encoded SHAP PNG per model.
            ai_insights:     AI-generated narrative text (shown on final slide).

        Returns:
            The ``output_path`` after saving.
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(self._SLIDE_W)
        prs.slide_height = Inches(self._SLIDE_H)

        # ----------------------------------------------------------------
        # Slide 1 – Title
        # ----------------------------------------------------------------
        self._add_title_slide(prs, title, task_type)

        # ----------------------------------------------------------------
        # Slide 2 – Dataset Overview
        # ----------------------------------------------------------------
        self._add_dataset_slide(prs, dataset_summary)

        # ----------------------------------------------------------------
        # Slide 3 – Model Comparison Table
        # ----------------------------------------------------------------
        self._add_model_comparison_slide(prs, metrics_table)

        # ----------------------------------------------------------------
        # Slide 4 – Metrics Summary (top-5 models bar chart as text)
        # ----------------------------------------------------------------
        self._add_metrics_summary_slide(prs, metrics_table, task_type)

        # ----------------------------------------------------------------
        # Slides 5 … N – SHAP plots
        # ----------------------------------------------------------------
        for model_name, b64_png in shap_images.items():
            self._add_shap_slide(prs, model_name, b64_png)

        # ----------------------------------------------------------------
        # Final Slide – Leaderboard + AI insights + Recommendations
        # ----------------------------------------------------------------
        self._add_final_slide(prs, leaderboard, ai_insights)

        prs.save(str(output_path))
        logger.info("PPTX report written to %s", output_path)
        return output_path

    # ------------------------------------------------------------------
    # Individual slide builders
    # ------------------------------------------------------------------

    def _add_title_slide(self, prs: Presentation, title: str, task_type: str) -> None:
        """Slide 1: Full-bleed title with CodeAlpha branding."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])

        # Background
        _add_rect(slide, 0, 0, self._SLIDE_W, self._SLIDE_H, _DARK)
        # Accent bar
        _add_rect(slide, 0, self._SLIDE_H * 0.62, self._SLIDE_W, 0.06, _ACCENT)

        # Company name
        _add_textbox(
            slide, "CodeAlpha",
            left=0.5, top=1.2, width=9.0, height=0.9,
            font_size=44, bold=True, color=_WHITE,
            align=PP_ALIGN.CENTER,
        )
        # Tagline
        _add_textbox(
            slide, "Enterprise AI Platform",
            left=0.5, top=2.1, width=9.0, height=0.5,
            font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE),
            align=PP_ALIGN.CENTER,
        )
        # Report title
        _add_textbox(
            slide, title,
            left=0.5, top=3.0, width=9.0, height=0.8,
            font_size=22, bold=True, color=_ACCENT,
            align=PP_ALIGN.CENTER,
        )
        # Task type badge
        _add_textbox(
            slide, f"Task: {task_type.title()}  •  Automated ML Experiment Report",
            left=0.5, top=3.9, width=9.0, height=0.4,
            font_size=12, italic=True, color=RGBColor(0xA5, 0xB4, 0xFC),
            align=PP_ALIGN.CENTER,
        )
        # Footer
        _add_textbox(
            slide, "Generated by CodeAlpha AutoML Pipeline",
            left=0.5, top=self._SLIDE_H - 0.45, width=9.0, height=0.35,
            font_size=9, italic=True, color=RGBColor(0x6B, 0x72, 0x80),
            align=PP_ALIGN.CENTER,
        )

    def _add_dataset_slide(
        self, prs: Presentation, dataset_summary: dict[str, Any]
    ) -> None:
        """Slide 2: Dataset overview facts."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        self._add_slide_header(slide, "Dataset Overview", "2")

        items = [(k.replace("_", " ").title(), str(v)) for k, v in dataset_summary.items()]
        n = len(items)
        cols = 2
        rows = (n + cols - 1) // cols

        card_w = 4.2
        card_h = 0.65
        margin_l = 0.4
        margin_top = 1.3
        gap_x = 0.3
        gap_y = 0.15

        for i, (key, val) in enumerate(items):
            row, col = divmod(i, cols)
            x = margin_l + col * (card_w + gap_x)
            y = margin_top + row * (card_h + gap_y)
            if y + card_h > self._SLIDE_H - 0.5:
                break  # safety clip
            _add_rect(slide, x, y, card_w, card_h, _LIGHT, _PRIMARY, 1)
            _add_textbox(
                slide, key,
                left=x + 0.1, top=y + 0.05, width=1.5, height=0.25,
                font_size=8, bold=True, color=_PRIMARY,
            )
            _add_textbox(
                slide, val,
                left=x + 0.1, top=y + 0.30, width=card_w - 0.2, height=0.28,
                font_size=12, bold=True, color=_DARK,
            )

    def _add_model_comparison_slide(
        self, prs: Presentation, metrics_table: list[dict[str, Any]]
    ) -> None:
        """Slide 3: Model comparison table."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        self._add_slide_header(slide, "Model Comparison", "3")

        if not metrics_table:
            _add_textbox(slide, "No metrics data available.", 0.5, 2, 9, 1, color=_DARK)
            return

        keys = list(metrics_table[0].keys())
        max_cols = 7
        keys = keys[:max_cols]

        col_width = min(1.35, 9.0 / len(keys))
        table_left = (self._SLIDE_W - col_width * len(keys)) / 2
        table_top = 1.3
        row_h = 0.42

        # Header
        for j, key in enumerate(keys):
            x = table_left + j * col_width
            _add_rect(slide, x, table_top, col_width, row_h, _PRIMARY)
            _add_textbox(
                slide, key.replace("_", "\n").title(),
                left=x + 0.02, top=table_top + 0.02,
                width=col_width - 0.04, height=row_h - 0.04,
                font_size=7, bold=True, color=_WHITE, align=PP_ALIGN.CENTER,
            )

        max_rows = min(len(metrics_table), 8)
        for i, row_data in enumerate(metrics_table[:max_rows]):
            y = table_top + (i + 1) * row_h
            bg = _LIGHT if i % 2 == 0 else _WHITE
            for j, key in enumerate(keys):
                x = table_left + j * col_width
                _add_rect(slide, x, y, col_width, row_h, bg, _PRIMARY, 0)
                val = _format_val(row_data.get(key))
                _add_textbox(
                    slide, val,
                    left=x + 0.02, top=y + 0.08,
                    width=col_width - 0.04, height=row_h - 0.08,
                    font_size=7, color=_DARK, align=PP_ALIGN.CENTER,
                )

    def _add_metrics_summary_slide(
        self,
        prs: Presentation,
        metrics_table: list[dict[str, Any]],
        task_type: str,
    ) -> None:
        """Slide 4: Textual metrics overview with top-model highlights."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        self._add_slide_header(slide, "Performance Metrics Summary", "4")

        if not metrics_table:
            _add_textbox(slide, "No metrics data.", 0.5, 2, 9, 1, color=_DARK)
            return

        # Sort by accuracy descending (best first)
        key = "accuracy" if "accuracy" in metrics_table[0] else list(metrics_table[0].keys())[1]
        sorted_models = sorted(
            metrics_table,
            key=lambda r: float(r.get(key, 0) or 0),
            reverse=True,
        )

        top3 = sorted_models[:3]
        card_w = 2.8
        card_h = 2.8
        margin_l = 0.55
        gap = 0.25
        top_y = 1.4

        for i, model in enumerate(top3):
            x = margin_l + i * (card_w + gap)
            medal = ["🥇", "🥈", "🥉"][i]
            model_name = model.get("model_name", model.get("model", f"Model {i+1}"))
            _add_rect(slide, x, top_y, card_w, card_h, _LIGHT, _PRIMARY, 1)
            _add_textbox(
                slide, f"{medal} #{i+1}",
                left=x + 0.1, top=top_y + 0.1, width=card_w - 0.2, height=0.4,
                font_size=16, bold=True, color=_SECONDARY, align=PP_ALIGN.CENTER,
            )
            _add_textbox(
                slide, model_name,
                left=x + 0.1, top=top_y + 0.55, width=card_w - 0.2, height=0.4,
                font_size=11, bold=True, color=_DARK, align=PP_ALIGN.CENTER,
            )
            y_offset = 1.05
            metric_keys = [k for k in model.keys() if k not in {"model", "model_name"}]
            for mk in metric_keys[:5]:
                val = model.get(mk, "—")
                if isinstance(val, float):
                    val = f"{val:.4f}"
                _add_textbox(
                    slide, f"{mk.replace('_',' ').title()}: {val}",
                    left=x + 0.1, top=top_y + y_offset, width=card_w - 0.2, height=0.3,
                    font_size=8, color=_DARK, align=PP_ALIGN.LEFT,
                )
                y_offset += 0.32

    def _add_shap_slide(
        self, prs: Presentation, model_name: str, b64_png: str
    ) -> None:
        """One slide per model: SHAP summary plot."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        self._add_slide_header(slide, f"SHAP Explainability — {model_name}", "")

        try:
            raw = base64.b64decode(b64_png)
            img_stream = BytesIO(raw)
            img_left = Inches(1.0)
            img_top = Inches(1.3)
            img_w = Inches(8.0)
            img_h = Inches(5.5)
            slide.shapes.add_picture(img_stream, img_left, img_top, img_w, img_h)
        except Exception as exc:
            logger.warning("Could not embed SHAP image for '%s': %s", model_name, exc)
            _add_textbox(
                slide, "[SHAP image could not be loaded]",
                left=0.5, top=3.0, width=9.0, height=0.5,
                font_size=11, italic=True, color=RGBColor(0xDC, 0x26, 0x26),
                align=PP_ALIGN.CENTER,
            )

        _add_textbox(
            slide,
            "SHAP summary plot — positive values increase predictions, "
            "negative values decrease them",
            left=0.5, top=self._SLIDE_H - 0.45, width=9.0, height=0.35,
            font_size=8, italic=True, color=RGBColor(0x6B, 0x72, 0x80),
            align=PP_ALIGN.CENTER,
        )

    def _add_final_slide(
        self,
        prs: Presentation,
        leaderboard: list[dict[str, Any]],
        ai_insights: str,
    ) -> None:
        """Final slide: Leaderboard + AI insights snippet."""
        slide = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        self._add_slide_header(slide, "Leaderboard & Key Insights", "")

        # Leaderboard mini-table (top-5)
        top5 = leaderboard[:5]
        if top5:
            lb_keys = [k for k in top5[0].keys()][:5]
            col_w = 1.6
            row_h = 0.38
            tbl_left = 0.4
            tbl_top = 1.3

            # Header
            for j, key in enumerate(lb_keys):
                x = tbl_left + j * col_w
                _add_rect(slide, x, tbl_top, col_w, row_h, _PRIMARY)
                _add_textbox(
                    slide, key.replace("_", " ").title(),
                    left=x + 0.02, top=tbl_top + 0.04,
                    width=col_w - 0.04, height=row_h - 0.04,
                    font_size=7, bold=True, color=_WHITE, align=PP_ALIGN.CENTER,
                )

            for i, entry in enumerate(top5):
                y = tbl_top + (i + 1) * row_h
                bg = RGBColor(0xFE, 0xF9, 0xC3) if i == 0 else (
                    _LIGHT if i % 2 == 0 else _WHITE
                )
                for j, key in enumerate(lb_keys):
                    x = tbl_left + j * col_w
                    _add_rect(slide, x, y, col_w, row_h, bg)
                    val = _format_val(entry.get(key))
                    _add_textbox(
                        slide, val,
                        left=x + 0.02, top=y + 0.06,
                        width=col_w - 0.04, height=row_h - 0.06,
                        font_size=7, bold=(i == 0), color=_DARK, align=PP_ALIGN.CENTER,
                    )

        # AI Insights panel
        insight_lines = [
            ln.strip() for ln in ai_insights.split("\n")
            if ln.strip() and not ln.startswith("#")
        ]
        insight_snippet = " ".join(insight_lines)[:400]
        if insight_snippet:
            _add_rect(slide, 8.2, 1.3, 1.6, 5.5, _LIGHT, _SECONDARY, 1)  # background panel - hidden
            _add_textbox(
                slide, "🤖 AI Insights",
                left=8.22, top=1.35, width=1.55, height=0.3,
                font_size=8, bold=True, color=_SECONDARY, align=PP_ALIGN.LEFT,
            )
            _add_textbox(
                slide, insight_snippet,
                left=8.22, top=1.7, width=1.55, height=4.8,
                font_size=7, color=_DARK, align=PP_ALIGN.LEFT,
                word_wrap=True,
            )

        # Footer call-to-action
        _add_rect(slide, 0, self._SLIDE_H - 0.55, self._SLIDE_W, 0.55, _DARK)
        _add_textbox(
            slide,
            "Thank you  •  Download the full PDF / DOCX report from CodeAlpha Platform",
            left=0, top=self._SLIDE_H - 0.45, width=self._SLIDE_W, height=0.35,
            font_size=10, bold=True, color=_WHITE, align=PP_ALIGN.CENTER,
        )

    # ------------------------------------------------------------------
    # Shared slide header band
    # ------------------------------------------------------------------

    def _add_slide_header(self, slide, title: str, slide_num: str) -> None:
        """Draw the top header band with title and optional slide number."""
        _add_rect(slide, 0, 0, self._SLIDE_W, 1.1, _PRIMARY)
        _add_rect(slide, 0, 1.1, self._SLIDE_W, 0.05, _ACCENT)
        _add_textbox(
            slide, "CodeAlpha  |  Enterprise AI Platform",
            left=0.3, top=0.05, width=7.0, height=0.35,
            font_size=9, color=RGBColor(0xC7, 0xD2, 0xFE),
        )
        _add_textbox(
            slide, title,
            left=0.3, top=0.45, width=8.5, height=0.55,
            font_size=20, bold=True, color=_WHITE,
        )
        if slide_num:
            _add_textbox(
                slide, f"Slide {slide_num}",
                left=9.0, top=0.05, width=0.9, height=0.35,
                font_size=9, color=RGBColor(0xC7, 0xD2, 0xFE),
                align=PP_ALIGN.RIGHT,
            )
